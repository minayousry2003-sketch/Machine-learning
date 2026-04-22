from pptx import Presentation
from pptx.util import Inches
import os

# إنشاء العرض التقديمي
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# محتوى الشرائح الأساسية
slides_content = [
    {
        "title": "Diabetes Detection Web Application",
        "content": [
            "Machine Learning with Django",
            "• Combines data science and web development",
            "• Predicts diabetes risk using patient health data",
            "• Demonstrates real-world deployment of ML models through a web interface"
        ]
    },
    {
        "title": "Problem Description",
        "content": [
            "• Goal: Predict diabetes risk from patient health data",
            "• Importance: Diabetes is a major global health issue",
            "• Challenge: Many ML systems are inaccessible to non-technical users",
            "• Impact: Provides a simple, web-based tool for doctors, patients, and researchers"
        ]
    },
    {
        "title": "Dataset Overview",
        "content": [
            "• Source: diabetes.csv (768 rows, 9 input features + 1 outcome column)",
            "• Features: Pregnancies, Glucose, BloodPressure, SkinThickness, Insulin, BMI, DiabetesPedigreeFunction, Age",
            "• Target: Outcome (0 = Negative, 1 = Positive)",
            "• Observations: Missing values (zeros in Insulin, SkinThickness) → preprocessing required",
            "• Why this dataset? Widely used in academic research, balanced mix of positive/negative cases"
        ]
    },
    {
        "title": "Model Training",
        "content": [
            "• Algorithm: Logistic Regression",
            "• Reason: Simple, interpretable, effective for binary classification",
            "• Workflow:",
            "   - Split dataset into training/testing sets",
            "   - Train model on training data",
            "   - Evaluate predictions on test data",
            "• Code Snippet:",
            "  x = data.drop('Outcome', axis=1)",
            "  y = data['Outcome']",
            "  x_train, x_test, y_train, y_test = train_test_split(x, y, test_size=0.2)",
            "  model = LogisticRegression(max_iter=200)",
            "  model.fit(x_train, y_train)"
        ]
    },
    {
        "title": "Web Integration (Django Views)",
        "content": [
            "• Purpose: Connect ML model with web interface",
            "• Functions:",
            "   - home() → Homepage",
            "   - predict() → Input form",
            "   - result() → Runs ML model, returns prediction",
            "• Code Snippet:",
            "  def result(request):",
            "      val1 = float(request.GET.get('n1'))",
            "      ...",
            "      pred = model.predict([[val1, val2, ..., val8]])",
            "      result2 = \"Positive\" if pred[0] == 1 else \"Negative\"",
            "      return render(request, 'predict.html', {\"result2\": result2})"
        ]
    },
    {
        "title": "Evaluation Metrics",
        "content": [
            "• Accuracy: ~77–80% baseline",
            "• Precision & Recall: Measure true positives vs false negatives",
            "• F1-score: Balances precision and recall",
            "• Visuals: Confusion matrix, ROC curve, accuracy comparison chart",
            "• Interpretation: Logistic Regression is effective, but advanced models could improve accuracy"
        ]
    },
    {
        "title": "Frontend (Templates)",
        "content": [
            "• home.html: Landing page introducing project",
            "• predict.html: Input form + prediction result",
            "• Static images: Enhance user experience",
            "• Design Goal: Simple, clean, accessible for non-technical users"
        ]
    },
    {
        "title": "Expected Output",
        "content": [
            "• Workflow:",
            "   - User enters 8 health values",
            "   - System processes input through ML model",
            "   - Output displayed as Positive or Negative",
            "• Examples:",
            "   - Input: Glucose=148, BMI=33.6, Age=50 → Output: Positive",
            "   - Input: Glucose=85, BMI=26.6, Age=31 → Output: Negative"
        ]
    },
    {
        "title": "Conclusion",
        "content": [
            "• Summary: End-to-end solution combining ML, Django, and UI",
            "• Benefits: Accessible predictions for healthcare applications",
            "• Future Work:",
            "   - Save/reuse trained models instead of retraining each request",
            "   - Deploy on cloud (Heroku, AWS, Azure)",
            "   - Add more models for comparison",
            "   - Enhance UI with charts and explanations"
        ]
    }
]

# إضافة الشرائح الأساسية
for slide_data in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = slide_data["title"]
    content.text = "\n".join(slide_data["content"])

# إضافة شريحة Confusion Matrix
slide_confusion = prs.slides.add_slide(prs.slide_layouts[1])
slide_confusion.shapes.title.text = "Confusion Matrix"
slide_confusion.placeholders[1].text = (
    "Visual representation of model predictions:\n"
    "• True Positives (TP)\n"
    "• True Negatives (TN)\n"
    "• False Positives (FP)\n"
    "• False Negatives (FN)\n"
    "Helps evaluate classification performance."
)

# إضافة شريحة ROC Curve
slide_roc = prs.slides.add_slide(prs.slide_layouts[1])
slide_roc.shapes.title.text = "ROC Curve"
slide_roc.placeholders[1].text = (
    "Receiver Operating Characteristic (ROC) Curve:\n"
    "• Plots True Positive Rate vs False Positive Rate\n"
    "• AUC (Area Under Curve) indicates model performance\n"
    "• Higher AUC = better discrimination ability"
)

# حفظ الملف مباشرة في المجلد الحالي
output_path = "Diabetes_Detection_WebApp_Presentation.pptx"
prs.save(output_path)

print(f"Presentation saved as {output_path}")
